from pptx import Presentation

def extract_text_from_pptx(input_path: str) -> str:
    """
    Extracts text from a .pptx file (slide titles, bullet points, etc.) and returns it as a single string.
    """
    prs = Presentation(input_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
